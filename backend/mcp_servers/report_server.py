#!/usr/bin/env python3
"""
SVA Report Generation Server - PDF/PPT Creation
Generates downloadable reports from analysis data
"""
import os
import json
import logging
from datetime import datetime
from typing import Dict, Any, Optional, List

# PDF Generation
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# PPT Generation  
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

logger = logging.getLogger(__name__)

class ReportGenerationServer:
    def __init__(self):
        self.output_dir = os.path.join(os.path.dirname(__file__), "..", "generated_reports")
        self.chat_storage_dir = os.path.join(os.path.dirname(__file__), "..", "chat_storage")
        os.makedirs(self.output_dir, exist_ok=True)
        logger.info(f"📄 Report generation server initialized. Output directory: {self.output_dir}")
        
    def generate_report_from_chat(self, chat_id: str, video_filename: str, format_type: str = "pdf") -> Dict[str, Any]:
        """Generate report by reading actual chat data"""
        try:
            # Read the chat file
            chat_file = os.path.join(self.chat_storage_dir, f"{chat_id}.json")
            if not os.path.exists(chat_file):
                return {
                    "success": False,
                    "error": f"Chat file not found: {chat_id}",
                    "filepath": None
                }
            
            with open(chat_file, 'r') as f:
                chat_data = json.load(f)
            
            # Extract analysis data from chat messages
            analysis_data = self._extract_analysis_from_chat(chat_data)
            
            # Use the video filename extracted from chat data if available, otherwise use passed parameter
            actual_video_filename = analysis_data.get('video_filename', 'unknown')
            if actual_video_filename and actual_video_filename != 'unknown':
                video_filename = actual_video_filename
                logger.info(f"📽️ Using extracted video filename: {video_filename}")
            else:
                logger.info(f"📽️ Using passed video filename: {video_filename}")
            
            # Generate report with real data
            return self.generate_analysis_report(analysis_data, video_filename, format_type)
            
        except Exception as e:
            logger.error(f"❌ Failed to generate report from chat {chat_id}: {e}")
            return {
                "success": False,
                "error": f"Failed to read chat data: {str(e)}",
                "filepath": None
            }
    
    def _extract_analysis_from_chat(self, chat_data: Dict) -> Dict[str, Any]:
        """Extract real analysis data from SVA chat messages"""
        analysis_data = {
            'chat_id': chat_data.get('chat_id', 'unknown'),
            'video_filename': 'unknown',
            'messages': chat_data.get('messages', []),
            'spoken_text': [],
            'visual_text': [],
            'objects_detected': [],
            'key_topics': [],
            'executive_summary': 'Video Analysis Report',
            'technical_details': {},
            'analysis_timestamp': datetime.now().isoformat()
        }
        
        # Parse messages to extract analysis content
        for message in chat_data.get('messages', []):
            content = message.get('content', '')
            
            # Extract video filename from user messages
            if message.get('role') == 'user' and message.get('file_path'):
                file_path = message.get('file_path', '')
                user_content = message.get('content', '')
                
                # Method 1: Extract from file path
                if '/' in file_path:
                    path_filename = file_path.split('/')[-1]
                else:
                    path_filename = file_path
                
                # Method 2: Look for original filename in user message content
                # Users often mention the original filename in their message
                original_filename = None
                if user_content:
                    video_extensions = ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm']
                    words = user_content.split()
                    for word in words:
                        for ext in video_extensions:
                            if ext.lower() in word.lower():
                                # Found a word containing video extension
                                original_filename = word.strip('",.:;!?')
                                break
                        if original_filename:
                            break
                
                # Use original filename if found, otherwise use path filename
                if original_filename and original_filename != path_filename:
                    analysis_data['video_filename'] = original_filename
                    analysis_data['original_filename'] = original_filename
                    analysis_data['system_path'] = path_filename
                    logger.info(f"📽️ Found original filename: {original_filename} (system: {path_filename})")
                else:
                    analysis_data['video_filename'] = path_filename
                    logger.info(f"📽️ Using path filename: {path_filename}")
            
            # Extract spoken text
            if '🎤 **Spoken Text:**' in content:
                spoken_lines = []
                lines = content.split('\n')
                in_spoken_section = False
                for line in lines:
                    if '🎤 **Spoken Text:**' in line:
                        in_spoken_section = True
                        continue
                    elif line.startswith('👁️') or line.startswith('*Language:'):
                        in_spoken_section = False
                        continue
                    elif in_spoken_section and '"' in line:
                        # Extract quoted text from numbered lines like: 1. "text here"
                        if '"' in line:
                            start_quote = line.find('"')
                            end_quote = line.rfind('"')
                            if start_quote != -1 and end_quote != -1 and start_quote != end_quote:
                                spoken_text = line[start_quote+1:end_quote]
                                if spoken_text.strip():
                                    spoken_lines.append(spoken_text.strip())
                analysis_data['spoken_text'] = spoken_lines
            
            # Extract visual text
            if '👁️ **Visual Text (On-screen):**' in content:
                visual_items = []
                lines = content.split('\n')
                in_visual_section = False
                for line in lines:
                    if '👁️ **Visual Text (On-screen):**' in line:
                        in_visual_section = True
                        continue
                    elif line.startswith('🎯') or line.startswith('📊'):
                        in_visual_section = False
                        continue
                    elif in_visual_section and '"' in line and '%' in line:
                        # Extract visual text items
                        parts = line.split('"')
                        if len(parts) >= 2:
                            visual_text = parts[1]
                            if visual_text and visual_text not in visual_items:
                                visual_items.append(visual_text)
                analysis_data['visual_text'] = visual_items
            
            # Extract objects detected
            if 'Video shows:' in content:
                objects_line = content.split('Video shows:')[1].split('\n')[0].strip()
                objects = [obj.strip() for obj in objects_line.split(',')]
                analysis_data['objects_detected'] = objects
            
            # Extract key topics
            if '💡 **Key Messages:**' in content or '🎯 **Main Themes:**' in content:
                topics = []
                lines = content.split('\n')
                for line in lines:
                    if ('•' in line or '•' in line) and ('mentioned' in line or 'success' in line or 'comfort' in line):
                        topic = line.strip('• ').strip()
                        if topic:
                            topics.append(topic)
                analysis_data['key_topics'] = topics
        
        # Create executive summary from extracted data
        if analysis_data['spoken_text'] or analysis_data['visual_text'] or analysis_data['objects_detected']:
            summary_parts = []
            if analysis_data['spoken_text']:
                summary_parts.append(f"Transcribed {len(analysis_data['spoken_text'])} spoken segments")
            if analysis_data['visual_text']:
                summary_parts.append(f"detected {len(analysis_data['visual_text'])} on-screen text elements")
            if analysis_data['objects_detected']:
                summary_parts.append(f"identified objects: {', '.join(analysis_data['objects_detected'])}")
            
            analysis_data['executive_summary'] = f"Comprehensive video analysis completed. {' and '.join(summary_parts)}."
        
        # Technical details
        analysis_data['technical_details'] = {
            'total_messages': len(chat_data.get('messages', [])),
            'spoken_segments': len(analysis_data['spoken_text']),
            'visual_elements': len(analysis_data['visual_text']),
            'objects_count': len(analysis_data['objects_detected']),
            'chat_duration': self._calculate_chat_duration(chat_data.get('messages', []))
        }
        
        return analysis_data
    
    def _calculate_chat_duration(self, messages: List[Dict]) -> str:
        """Calculate duration of chat session"""
        if len(messages) < 2:
            return "< 1 minute"
        
        first_timestamp = messages[0].get('timestamp', 0)
        last_timestamp = messages[-1].get('timestamp', 0)
        duration_ms = last_timestamp - first_timestamp
        duration_minutes = duration_ms // 60000
        
        if duration_minutes < 1:
            return "< 1 minute"
        elif duration_minutes < 60:
            return f"{duration_minutes} minutes"
        else:
            hours = duration_minutes // 60
            minutes = duration_minutes % 60
            return f"{hours}h {minutes}m"
        
    def generate_analysis_report(self, analysis_data: Dict[str, Any], video_filename: str, format_type: str = "pdf") -> Dict[str, Any]:
        """Generate comprehensive report from SVA analysis data"""
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            clean_filename = self._clean_filename(video_filename)
            
            logger.info(f"📊 Generating {format_type.upper()} report for {video_filename}")
            
            if format_type.lower() == "pdf" and PDF_AVAILABLE:
                return self._generate_pdf_report(analysis_data, clean_filename, timestamp)
            elif format_type.lower() == "ppt" and PPT_AVAILABLE:
                return self._generate_ppt_report(analysis_data, clean_filename, timestamp)
            else:
                return self._generate_text_report(analysis_data, clean_filename, timestamp)
                
        except Exception as e:
            logger.error(f"❌ Report generation failed: {e}")
            return {
                "success": False,
                "error": f"Report generation failed: {str(e)}",
                "filepath": None
            }
    
    def _clean_filename(self, filename: str) -> str:
        """Clean filename for safe file operations"""
        import re
        # Remove file extension and clean up
        clean = re.sub(r'[^\w\-_\.]', '_', str(filename))
        if '.' in clean:
            clean = clean.split('.')[0]
        return clean[:50]  # Limit length
    
    def _generate_pdf_report(self, data: Dict[str, Any], video_name: str, timestamp: str) -> Dict[str, Any]:
        """Generate professional PDF report"""
        filename = f"SVA_Analysis_{video_name}_{timestamp}.pdf"
        filepath = os.path.join(self.output_dir, filename)
        
        # Create PDF document
        doc = SimpleDocTemplate(
            filepath, 
            pagesize=A4, 
            rightMargin=72, 
            leftMargin=72, 
            topMargin=72, 
            bottomMargin=72
        )
        
        # Define styles
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=22,
            textColor=colors.darkblue,
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        section_style = ParagraphStyle(
            'SectionHeader',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.darkblue,
            spaceBefore=20,
            spaceAfter=12,
            fontName='Helvetica-Bold'
        )
        
        subsection_style = ParagraphStyle(
            'SubsectionHeader',
            parent=styles['Heading3'],
            fontSize=12,
            textColor=colors.darkslategray,
            spaceBefore=15,
            spaceAfter=8,
            fontName='Helvetica-Bold'
        )
        
        # Build content
        story = []
        
        # Title page
        story.append(Paragraph("🎬 SVA Video Analysis Report", title_style))
        story.append(Spacer(1, 20))
        
        # Metadata table
        metadata_data = [
            ['📁 Video File:', video_name],
            ['📅 Analysis Date:', datetime.now().strftime('%B %d, %Y at %H:%M')],
            ['🤖 Analysis Engine:', 'SVA HuggingFace Compliant System'],
            ['📊 Report Format:', 'Comprehensive Analysis']
        ]
        
        metadata_table = Table(metadata_data, colWidths=[2*inch, 4*inch])
        metadata_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.lightblue),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ]))
        
        story.append(metadata_table)
        story.append(Spacer(1, 30))
        
        # Executive Summary
        story.append(Paragraph("📋 Executive Summary", section_style))
        summary_text = self._generate_executive_summary(data)
        story.append(Paragraph(summary_text, styles['Normal']))
        story.append(Spacer(1, 20))
        
        # Speech Transcription Section
        if 'spoken_text' in data and data['spoken_text']:
            story.append(Paragraph("🎤 Speech Transcription Analysis", section_style))
            
            spoken_data = data['spoken_text']
            language = data.get('language', 'Auto-detected')
            
            story.append(Paragraph(f"<b>Language Detected:</b> {language}", subsection_style))
            
            if isinstance(spoken_data, list):
                story.append(Paragraph("<b>Transcribed Segments:</b>", subsection_style))
                for i, text_chunk in enumerate(spoken_data[:8], 1):  # Limit to 8 segments
                    story.append(Paragraph(f"{i}. \"{text_chunk}\"", styles['Normal']))
                    story.append(Spacer(1, 6))
            else:
                story.append(Paragraph(f"<b>Full Transcription:</b>", subsection_style))
                story.append(Paragraph(f"\"{spoken_data}\"", styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Visual Text Section
        if 'visual_text' in data and data['visual_text']:
            story.append(Paragraph("👁️ Visual Text Recognition", section_style))
            visual_data = data['visual_text']
            
            if isinstance(visual_data, list) and visual_data:
                # Create table for visual text
                visual_table_data = [['Text Element', 'Confidence', 'Timestamp']]
                for item in visual_data[:15]:  # Limit to 15 items
                    if isinstance(item, dict):
                        text_val = item.get('text', 'N/A')
                        confidence = f"{item.get('confidence', 0):.1%}" if 'confidence' in item else 'N/A'
                        timestamp = f"{item.get('timestamp', 'N/A')}"
                        visual_table_data.append([text_val, confidence, timestamp])
                
                visual_table = Table(visual_table_data, colWidths=[3*inch, 1.5*inch, 1.5*inch])
                visual_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(visual_table)
            story.append(Spacer(1, 20))
        
        # Object Detection Section
        if 'objects' in data and data['objects']:
            story.append(Paragraph("🎯 Object Detection Results", section_style))
            objects_data = data['objects']
            
            if isinstance(objects_data, list) and objects_data:
                # Create table for objects
                objects_table_data = [['Object Type', 'Instances', 'Avg Confidence']]
                for obj in objects_data[:12]:  # Limit to 12 objects
                    if isinstance(obj, dict):
                        label = obj.get('label', 'Unknown')
                        count = str(obj.get('count', 1))
                        confidence = f"{obj.get('confidence', 0):.1%}" if 'confidence' in obj else 'N/A'
                        objects_table_data.append([label, count, confidence])
                
                objects_table = Table(objects_table_data, colWidths=[2.5*inch, 1.5*inch, 1.5*inch])
                objects_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.darkblue),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.lightcyan),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(objects_table)
            story.append(Spacer(1, 20))
        
        # Topic Analysis Section
        if 'topics' in data and data['topics']:
            story.append(Paragraph("📊 Intelligent Content Analysis", section_style))
            topics_data = data['topics']
            
            if 'themes' in topics_data and topics_data['themes']:
                story.append(Paragraph("<b>🎯 Main Themes Identified:</b>", subsection_style))
                for i, theme in enumerate(topics_data['themes'][:5], 1):
                    story.append(Paragraph(f"{i}. {theme}", styles['Normal']))
                story.append(Spacer(1, 10))
            
            if 'key_phrases' in topics_data and topics_data['key_phrases']:
                story.append(Paragraph("<b>💡 Key Messages:</b>", subsection_style))
                for phrase in topics_data['key_phrases'][:3]:
                    story.append(Paragraph(f"• \"{phrase}\"", styles['Normal']))
                story.append(Spacer(1, 10))
                
            if 'content_type' in topics_data:
                story.append(Paragraph(f"<b>📖 Content Classification:</b> {topics_data['content_type']}", subsection_style))
                
            if 'setting' in topics_data:
                story.append(Paragraph(f"<b>🎬 Scene Setting:</b> {topics_data['setting']}", subsection_style))
            
            story.append(Spacer(1, 20))
        
        # Technical Details
        story.append(Paragraph("⚙️ Technical Analysis Details", section_style))
        tech_details = [
            "• Audio Analysis: HuggingFace Whisper with automatic language detection",
            "• Visual Recognition: DETR object detection + TrOCR text recognition",
            "• Content Analysis: Advanced NLP with thematic extraction",
            "• Processing: Local AI models (fully offline)",
            f"• Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        ]
        for detail in tech_details:
            story.append(Paragraph(detail, styles['Normal']))
        
        # Build PDF
        doc.build(story)
        
        logger.info(f"✅ PDF report generated: {filename}")
        return {
            "success": True,
            "filepath": filepath,
            "filename": filename,
            "format": "pdf",
            "size": os.path.getsize(filepath) if os.path.exists(filepath) else 0
        }
    
    def _generate_ppt_report(self, data: Dict[str, Any], video_name: str, timestamp: str) -> Dict[str, Any]:
        """Generate PowerPoint presentation report"""
        filename = f"SVA_Analysis_{video_name}_{timestamp}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "🎬 SVA Video Analysis Report"
        subtitle.text = f"Analysis of: {video_name}\nGenerated: {datetime.now().strftime('%B %d, %Y at %H:%M')}\nPowered by SVA HuggingFace System"
        
        # Overview slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "📋 Analysis Overview"
        tf = body_shape.text_frame
        tf.text = f"📁 Video File: {video_name}"
        
        p = tf.add_paragraph()
        p.text = f"📅 Analysis Date: {datetime.now().strftime('%B %d, %Y at %H:%M')}"
        
        p = tf.add_paragraph()
        p.text = "🤖 Analysis Engine: SVA HuggingFace Compliant System"
        
        p = tf.add_paragraph()
        p.text = "🎯 Analysis Type: Comprehensive Video Content Analysis"
        
        # Executive Summary slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        title_shape.text = "📊 Executive Summary"
        
        tf = body_shape.text_frame
        summary = self._generate_executive_summary(data)
        tf.text = summary[:500] + "..." if len(summary) > 500 else summary
        
        # Speech Transcription slide
        if 'spoken_text' in data and data['spoken_text']:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.shapes.placeholders[1]
            title_shape.text = "🎤 Speech Transcription"
            
            tf = body_shape.text_frame
            spoken_data = data['spoken_text']
            language = data.get('language', 'Auto-detected')
            tf.text = f"Language: {language}"
            
            if isinstance(spoken_data, list):
                for text_chunk in spoken_data[:4]:  # Limit to 4 chunks for readability
                    p = tf.add_paragraph()
                    p.text = f"• \"{text_chunk[:100]}{'...' if len(text_chunk) > 100 else ''}\""
            else:
                p = tf.add_paragraph()
                p.text = f"• \"{str(spoken_data)[:200]}{'...' if len(str(spoken_data)) > 200 else ''}\""
        
        # Visual Text slide
        if 'visual_text' in data and data['visual_text']:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.shapes.placeholders[1]
            title_shape.text = "👁️ Visual Text Recognition"
            
            tf = body_shape.text_frame
            visual_data = data['visual_text']
            tf.text = f"Detected {len(visual_data)} visual text elements"
            
            if isinstance(visual_data, list):
                for item in visual_data[:6]:  # Limit to 6 items
                    if isinstance(item, dict):
                        text_val = item.get('text', 'N/A')
                        confidence = item.get('confidence', 0)
                        p = tf.add_paragraph()
                        p.text = f"• \"{text_val}\" ({confidence:.1%})"
        
        # Object Detection slide
        if 'objects' in data and data['objects']:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.shapes.placeholders[1]
            title_shape.text = "🎯 Object Detection"
            
            tf = body_shape.text_frame
            objects_data = data['objects']
            tf.text = f"Detected {len(objects_data)} object types"
            
            if isinstance(objects_data, list):
                for obj in objects_data[:8]:  # Limit to 8 objects
                    if isinstance(obj, dict):
                        label = obj.get('label', 'Unknown')
                        count = obj.get('count', 1)
                        confidence = obj.get('confidence', 0)
                        p = tf.add_paragraph()
                        p.text = f"• {label}: {count} instances ({confidence:.1%})"
        
        # Topic Analysis slide
        if 'topics' in data and data['topics']:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.shapes.placeholders[1]
            title_shape.text = "📊 Content Analysis"
            
            tf = body_shape.text_frame
            topics_data = data['topics']
            
            if 'content_type' in topics_data:
                tf.text = f"Content Type: {topics_data['content_type']}"
            
            if 'themes' in topics_data and topics_data['themes']:
                p = tf.add_paragraph()
                p.text = "Main Themes:"
                for theme in topics_data['themes'][:4]:
                    p = tf.add_paragraph()
                    p.text = f"• {theme}"
        
        # Conclusion slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        title_shape.text = "✅ Analysis Complete"
        
        tf = body_shape.text_frame
        tf.text = "Comprehensive video analysis completed successfully"
        
        p = tf.add_paragraph()
        p.text = f"Report generated on {datetime.now().strftime('%B %d, %Y at %H:%M')}"
        
        p = tf.add_paragraph()
        p.text = "Thank you for using SVA - Smart Video Analyzer"
        
        # Save presentation
        prs.save(filepath)
        
        logger.info(f"✅ PPT report generated: {filename}")
        return {
            "success": True,
            "filepath": filepath,
            "filename": filename,
            "format": "pptx",
            "size": os.path.getsize(filepath) if os.path.exists(filepath) else 0
        }
    
    def _generate_text_report(self, data: Dict[str, Any], video_name: str, timestamp: str) -> Dict[str, Any]:
        """Generate fallback text report"""
        filename = f"SVA_Analysis_{video_name}_{timestamp}.txt"
        filepath = os.path.join(self.output_dir, filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write("🎬 SVA VIDEO ANALYSIS REPORT\n")
            f.write("=" * 50 + "\n\n")
            f.write(f"Video File: {video_name}\n")
            
            # Show additional filename info if available
            if 'original_filename' in data and data['original_filename'] != video_name:
                f.write(f"Original Upload: {data['original_filename']}\n")
            if 'system_path' in data and data['system_path'] != video_name:
                f.write(f"System Path: {data['system_path']}\n")
                
            f.write(f"Analysis Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Analysis Engine: SVA HuggingFace Compliant System\n\n")
            
            # Executive Summary
            f.write("EXECUTIVE SUMMARY:\n")
            f.write("-" * 20 + "\n")
            f.write(self._generate_executive_summary(data) + "\n\n")
            
            # Speech Transcription
            if 'spoken_text' in data and data['spoken_text']:
                f.write("SPEECH TRANSCRIPTION:\n")
                f.write("-" * 20 + "\n")
                spoken_data = data['spoken_text']
                language = data.get('language', 'Auto-detected')
                f.write(f"Language: {language}\n\n")
                
                if isinstance(spoken_data, list):
                    for i, text_chunk in enumerate(spoken_data, 1):
                        f.write(f"{i}. \"{text_chunk}\"\n")
                else:
                    f.write(f"\"{spoken_data}\"\n")
                f.write("\n")
            
            # Visual Text
            if 'visual_text' in data and data['visual_text']:
                f.write("VISUAL TEXT RECOGNITION:\n")
                f.write("-" * 20 + "\n")
                visual_data = data['visual_text']
                for item in visual_data:
                    if isinstance(item, dict):
                        text_val = item.get('text', 'N/A')
                        confidence = item.get('confidence', 0)
                        timestamp = item.get('timestamp', 'N/A')
                        f.write(f"• \"{text_val}\" (Confidence: {confidence:.1%}, Time: {timestamp})\n")
                f.write("\n")
            
            # Objects
            if 'objects' in data and data['objects']:
                f.write("OBJECT DETECTION:\n")
                f.write("-" * 20 + "\n")
                objects_data = data['objects']
                for obj in objects_data:
                    if isinstance(obj, dict):
                        label = obj.get('label', 'Unknown')
                        count = obj.get('count', 1)
                        confidence = obj.get('confidence', 0)
                        f.write(f"• {label}: {count} instances (Confidence: {confidence:.1%})\n")
                f.write("\n")
            
            # Topics
            if 'topics' in data and data['topics']:
                f.write("CONTENT ANALYSIS:\n")
                f.write("-" * 20 + "\n")
                topics_data = data['topics']
                
                if 'content_type' in topics_data:
                    f.write(f"Content Type: {topics_data['content_type']}\n")
                
                if 'themes' in topics_data:
                    f.write("Main Themes:\n")
                    for theme in topics_data['themes']:
                        f.write(f"• {theme}\n")
                
                if 'key_phrases' in topics_data:
                    f.write("Key Messages:\n")
                    for phrase in topics_data['key_phrases']:
                        f.write(f"• \"{phrase}\"\n")
                f.write("\n")
            
            f.write("TECHNICAL DETAILS:\n")
            f.write("-" * 20 + "\n")
            f.write("• Audio Analysis: HuggingFace Whisper with automatic language detection\n")
            f.write("• Visual Recognition: DETR object detection + TrOCR text recognition\n")
            f.write("• Content Analysis: Advanced NLP with thematic extraction\n")
            f.write("• Processing: Local AI models (fully offline)\n")
            f.write(f"• Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        
        logger.info(f"✅ Text report generated: {filename}")
        return {
            "success": True,
            "filepath": filepath,
            "filename": filename,
            "format": "txt",
            "size": os.path.getsize(filepath)
        }
    
    def _generate_executive_summary(self, data: Dict[str, Any]) -> str:
        """Generate executive summary from analysis data"""
        summary_parts = []
        
        # Count total elements
        spoken_count = 0
        visual_count = 0
        object_count = 0
        
        if 'spoken_text' in data and data['spoken_text']:
            spoken_data = data['spoken_text']
            if isinstance(spoken_data, list):
                spoken_count = len(spoken_data)
            else:
                spoken_count = 1
        
        if 'visual_text' in data and data['visual_text']:
            visual_count = len(data['visual_text'])
        
        if 'objects' in data and data['objects']:
            object_count = len(data['objects'])
        
        # Build summary
        summary_parts.append(f"This comprehensive analysis identified {spoken_count + visual_count} text sources and {object_count} object types.")
        
        if spoken_count > 0:
            language = data.get('language', 'auto-detected language')
            summary_parts.append(f"Audio analysis detected {spoken_count} speech segments in {language}.")
        
        if visual_count > 0:
            summary_parts.append(f"Visual text recognition found {visual_count} on-screen text elements.")
        
        if object_count > 0:
            summary_parts.append(f"Object detection identified {object_count} distinct object types in the video.")
        
        if 'topics' in data and data['topics']:
            topics_data = data['topics']
            if 'content_type' in topics_data:
                summary_parts.append(f"Content analysis classified this as: {topics_data['content_type']}.")
        
        summary_parts.append("All analysis was performed using local AI models ensuring complete privacy and offline operation.")
        
        return " ".join(summary_parts)
    
    def list_generated_reports(self) -> Dict[str, Any]:
        """List all generated reports in the output directory"""
        try:
            reports = []
            if os.path.exists(self.output_dir):
                for filename in os.listdir(self.output_dir):
                    filepath = os.path.join(self.output_dir, filename)
                    if os.path.isfile(filepath):
                        reports.append({
                            "filename": filename,
                            "filepath": filepath,
                            "size": os.path.getsize(filepath),
                            "created": datetime.fromtimestamp(os.path.getctime(filepath)).strftime('%Y-%m-%d %H:%M:%S')
                        })
            
            return {
                "success": True,
                "reports": sorted(reports, key=lambda x: x['created'], reverse=True)
            }
            
        except Exception as e:
            logger.error(f"❌ Failed to list reports: {e}")
            return {
                "success": False,
                "error": str(e),
                "reports": []
            }
    
    def cleanup_old_reports(self, days_old: int = 7) -> Dict[str, Any]:
        """Clean up reports older than specified days"""
        try:
            deleted_count = 0
            cutoff_time = datetime.now().timestamp() - (days_old * 24 * 60 * 60)
            
            if os.path.exists(self.output_dir):
                for filename in os.listdir(self.output_dir):
                    filepath = os.path.join(self.output_dir, filename)
                    if os.path.isfile(filepath) and os.path.getctime(filepath) < cutoff_time:
                        os.remove(filepath)
                        deleted_count += 1
            
            return {
                "success": True,
                "deleted_count": deleted_count,
                "message": f"Cleaned up {deleted_count} reports older than {days_old} days"
            }
            
        except Exception as e:
            logger.error(f"❌ Failed to cleanup reports: {e}")
            return {
                "success": False,
                "error": str(e),
                "deleted_count": 0
            }

if __name__ == "__main__":
    # Test the report generation
    test_data = {
        "spoken_text": [
            "so you got to get comfortable being uncomfortable if you ever want to be successful",
            "if you stay in your comfort zone that's where you will fail"
        ],
        "language": "English",
        "visual_text": [
            {"text": "AD", "confidence": 0.85, "timestamp": 5.3},
            {"text": "1", "confidence": 0.85, "timestamp": 2.7}
        ],
        "objects": [
            {"label": "person", "count": 11, "confidence": 1.0},
            {"label": "tie", "count": 11, "confidence": 1.0}
        ],
        "topics": {
            "themes": ["Personal growth and stepping outside comfort zones", "Success mindset"],
            "content_type": "Motivational/Inspirational speech",
            "setting": "Professional office environment"
        }
    }
    
    server = ReportGenerationServer()
    result = server.generate_analysis_report(test_data, "sample_video", "pdf")
    print(f"Test result: {result}")