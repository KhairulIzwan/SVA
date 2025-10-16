"""
Generation MCP Server - Document generation component for SVA project
Implements Model Context Protocol for creating PDF reports and PowerPoint presentations
"""

import asyncio
import json
import logging
from typing import Dict, Any, Optional, List
from datetime import datetime
from pathlib import Path
import io
import base64

# Document generation imports
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER, TA_LEFT
except ImportError:
    print("⚠️ ReportLab not installed. PDF generation will be limited.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("⚠️ python-pptx not installed. PowerPoint generation will be limited.")

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class GenerationMCPServer:
    """MCP Server for document generation (PDF reports and PowerPoint presentations)"""
    
    def __init__(self):
        self.server_name = "generation"
        self.capabilities = [
            "generate_pdf_report",
            "create_powerpoint",
            "generate_summary_report",
            "create_visual_summary",
            "export_transcription_report",
            "create_multi_modal_report",
            "get_report_templates"
        ]
        self.supported_formats = [".pdf", ".pptx", ".docx"]
        
        # Report templates
        self.report_templates = {
            "basic_analysis": {
                "name": "Basic Video Analysis Report",
                "sections": ["Summary", "Transcription", "Visual Analysis", "Key Findings"]
            },
            "comprehensive": {
                "name": "Comprehensive Multi-Modal Report", 
                "sections": ["Executive Summary", "Audio Analysis", "Visual Content", "Objects Detected", "Text Extracted", "Recommendations"]
            },
            "transcription_only": {
                "name": "Transcription Report",
                "sections": ["Video Information", "Full Transcription", "Language Analysis", "Key Timestamps"]
            },
            "visual_only": {
                "name": "Visual Content Analysis",
                "sections": ["Scene Overview", "Objects Detected", "Text Content", "Visual Statistics"]
            }
        }
        
    async def initialize(self):
        """Initialize generation server"""
        logger.info("📄 Initializing Generation MCP Server...")
        
        try:
            # Check available libraries
            available_features = []
            
            try:
                import reportlab
                available_features.append("PDF generation")
            except ImportError:
                logger.warning("ReportLab not available - installing...")
                import subprocess
                subprocess.run(["pip", "install", "reportlab"], check=True)
                available_features.append("PDF generation")
            
            try:
                import pptx
                available_features.append("PowerPoint generation")
            except ImportError:
                logger.warning("python-pptx not available - installing...")
                import subprocess
                subprocess.run(["pip", "install", "python-pptx"], check=True)
                available_features.append("PowerPoint generation")
            
            logger.info("✅ Generation MCP server ready")
            return {
                "status": "success",
                "available_features": available_features
            }
            
        except Exception as e:
            logger.error(f"❌ Generation server initialization failed: {e}")
            return {"status": "error", "error": str(e)}
    
    async def generate_report(self, summary_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate a comprehensive analysis report from video analysis results"""
        try:
            logger.info("📊 Generating comprehensive video analysis report...")
            
            # Extract data from analysis results
            video_path = summary_data.get('video_path', 'Unknown')
            transcription_data = summary_data.get('transcription')
            vision_data = summary_data.get('vision')
            
            # Create report content
            report_sections = []
            
            # Video Information Section
            report_sections.append({
                "title": "📹 Video Analysis Summary",
                "content": f"File: {Path(video_path).name}\nAnalysis completed: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            })
            
            # Transcription Section
            if transcription_data and hasattr(transcription_data, 'text') and transcription_data.text:
                language = getattr(transcription_data, 'language', 'unknown')
                confidence = getattr(transcription_data, 'confidence', 0.0)
                report_sections.append({
                    "title": f"🎤 Audio Transcription ({language.upper()})",
                    "content": f"Confidence: {confidence:.2f}\n\nTranscribed Text:\n\"{transcription_data.text}\""
                })
            
            # Vision Analysis Section
            if vision_data:
                objects_count = len(getattr(vision_data, 'objects_detected', []))
                texts_count = len(getattr(vision_data, 'text_extracted', []))
                
                vision_summary = f"Objects detected: {objects_count}\nText elements found: {texts_count}"
                vision_summary += f"\nFrames analyzed: {getattr(vision_data, 'frames_analyzed', 0)}"
                vision_summary += f"\nProcessing method: {getattr(vision_data, 'processing_method', 'unknown')}"
                vision_summary += f"\nScene description: {getattr(vision_data, 'scene_description', 'N/A')}"
                
                # Add detailed object information
                if hasattr(vision_data, 'objects_detected') and vision_data.objects_detected:
                    vision_summary += f"\n\n📦 Detected Objects (Top {min(10, len(vision_data.objects_detected))}):"
                    for i, obj in enumerate(vision_data.objects_detected[:10]):
                        confidence = getattr(obj, 'confidence', 0.0)
                        frame = getattr(obj, 'frame', 0)
                        timestamp = getattr(obj, 'timestamp', 0.0)
                        vision_summary += f"\n  {i+1}. {obj.class_name} (confidence: {confidence:.2f}, frame: {frame}, time: {timestamp:.1f}s)"
                
                # Add detailed text extraction results  
                if hasattr(vision_data, 'text_extracted') and vision_data.text_extracted:
                    vision_summary += f"\n\n📖 Extracted Text (Top {min(10, len(vision_data.text_extracted))}):"
                    for i, text in enumerate(vision_data.text_extracted[:10]):
                        confidence = getattr(text, 'confidence', 0.0)
                        frame = getattr(text, 'frame', 0)
                        timestamp = getattr(text, 'timestamp', 0.0)
                        text_content = text.text[:100] + "..." if len(text.text) > 100 else text.text
                        vision_summary += f"\n  {i+1}. \"{text_content}\" (confidence: {confidence:.2f}, frame: {frame}, time: {timestamp:.1f}s)"
                
                report_sections.append({
                    "title": "👁️ Visual Content Analysis",
                    "content": vision_summary
                })
            
            # Generate final report text
            report_text = "\n\n".join([f"{section['title']}\n{'='*50}\n{section['content']}" 
                                     for section in report_sections])
            
            # Add summary statistics
            total_analysis_types = len([s for s in [transcription_data, vision_data] if s])
            report_text += f"\n\n📊 Analysis Summary\n{'='*50}\n"
            report_text += f"Analysis types completed: {total_analysis_types}\n"
            report_text += f"Report generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            report_text += "Processing: HuggingFace compliant models"
            
            return {
                "status": "success",
                "report_type": "comprehensive_analysis",
                "report_text": report_text,
                "sections_generated": len(report_sections),
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"❌ Report generation failed: {e}")
            return {
                "status": "error",
                "error": str(e),
                "report_text": f"Report generation failed: {str(e)}"
            }
    
    async def process_request(self, request: Dict[str, Any]) -> Dict[str, Any]:
        """Main generation request handler"""
        action = request.get("action")
        request_id = request.get("request_id", f"req_{datetime.now().isoformat()}")
        
        logger.info(f"📄 Processing generation request {request_id}: {action}")
        
        try:
            if action == "generate_pdf_report":
                return await self._generate_pdf_report(request)
            elif action == "create_powerpoint":
                return await self._create_powerpoint(request)
            elif action == "generate_summary_report":
                return await self._generate_summary_report(request)
            elif action == "create_multi_modal_report":
                return await self._create_multi_modal_report(request)
            elif action == "export_transcription_report":
                return await self._export_transcription_report(request)
            elif action == "get_report_templates":
                return await self._get_report_templates()
            elif action == "get_capabilities":
                return await self._get_capabilities()
            else:
                return {
                    "request_id": request_id,
                    "status": "error",
                    "error": f"Unknown action: {action}",
                    "available_actions": self.capabilities
                }
                
        except Exception as e:
            logger.error(f"❌ Generation request processing failed: {e}")
            return {
                "request_id": request_id,
                "status": "error",
                "error": str(e)
            }
    
    async def _generate_pdf_report(self, request: Dict[str, Any]) -> Dict[str, Any]:
        """Generate comprehensive PDF report"""
        request_id = request.get("request_id", "unknown")
        analysis_data = request.get("analysis_data", {})
        template = request.get("template", "comprehensive")
        output_path = request.get("output_path", f"reports/video_analysis_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf")
        
        try:
            # Ensure reports directory exists
            Path("reports").mkdir(exist_ok=True)
            
            # Create PDF document
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
                alignment=TA_CENTER,
                textColor=colors.darkblue
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=14,
                spaceAfter=12,
                textColor=colors.darkblue
            )
            
            # Title
            story.append(Paragraph("SVA Video Analysis Report", title_style))
            story.append(Spacer(1, 12))
            
            # Metadata
            story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
            story.append(Paragraph(f"Video: {analysis_data.get('video_name', 'Unknown')}", styles['Normal']))
            story.append(Spacer(1, 20))
            
            # Executive Summary
            story.append(Paragraph("Executive Summary", heading_style))
            summary_text = self._generate_executive_summary(analysis_data)
            story.append(Paragraph(summary_text, styles['Normal']))
            story.append(Spacer(1, 15))
            
            # Audio Analysis Section
            if analysis_data.get("transcription"):
                story.append(Paragraph("Audio Analysis", heading_style))
                
                trans_data = analysis_data["transcription"]
                story.append(Paragraph(f"<b>Language Detected:</b> {trans_data.get('language', 'Unknown')}", styles['Normal']))
                story.append(Paragraph(f"<b>Confidence Score:</b> {trans_data.get('confidence_score', 0):.2f}", styles['Normal']))
                story.append(Paragraph(f"<b>Duration:</b> {trans_data.get('duration', 0):.2f} seconds", styles['Normal']))
                story.append(Spacer(1, 10))
                
                story.append(Paragraph("<b>Full Transcription:</b>", styles['Normal']))
                transcription_text = trans_data.get('text', 'No transcription available')
                story.append(Paragraph(transcription_text, styles['Normal']))
                story.append(Spacer(1, 15))
            
            # Visual Analysis Section
            if analysis_data.get("vision"):
                story.append(Paragraph("Visual Content Analysis", heading_style))
                
                vision_data = analysis_data["vision"]
                
                # Objects detected
                if vision_data.get("objects_detected"):
                    story.append(Paragraph("<b>Objects Detected:</b>", styles['Normal']))
                    
                    # Create table for objects
                    object_data = [["Object", "Confidence", "Count"]]
                    object_counts = {}
                    
                    for obj in vision_data["objects_detected"]:
                        class_name = obj["class_name"]
                        confidence = obj["confidence"]
                        object_counts[class_name] = object_counts.get(class_name, 0) + 1
                    
                    for obj_name, count in object_counts.items():
                        avg_confidence = sum(obj["confidence"] for obj in vision_data["objects_detected"] 
                                           if obj["class_name"] == obj_name) / count
                        object_data.append([obj_name.title(), f"{avg_confidence:.2f}", str(count)])
                    
                    object_table = Table(object_data)
                    object_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 14),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ]))
                    story.append(object_table)
                    story.append(Spacer(1, 15))
                
                # Text extracted
                if vision_data.get("text_found"):
                    story.append(Paragraph("<b>Text Content Extracted:</b>", styles['Normal']))
                    text_content = "; ".join([item["text"] for item in vision_data["text_found"]])
                    story.append(Paragraph(text_content, styles['Normal']))
                    story.append(Spacer(1, 15))
            
            # Key Findings
            story.append(Paragraph("Key Findings & Recommendations", heading_style))
            findings = self._generate_key_findings(analysis_data)
            for finding in findings:
                story.append(Paragraph(f"• {finding}", styles['Normal']))
            story.append(Spacer(1, 15))
            
            # Build PDF
            doc.build(story)
            
            # Get file size
            file_size = Path(output_path).stat().st_size
            
            logger.info(f"📄 PDF report generated: {output_path}")
            
            return {
                "request_id": request_id,
                "status": "success",
                "data": {
                    "output_path": output_path,
                    "file_size": file_size,
                    "format": "PDF",
                    "template_used": template,
                    "sections_included": ["Executive Summary", "Audio Analysis", "Visual Analysis", "Key Findings"]
                }
            }
            
        except Exception as e:
            logger.error(f"❌ PDF generation failed: {e}")
            return {
                "request_id": request_id,
                "status": "error",
                "error": str(e)
            }
    
    async def _create_powerpoint(self, request: Dict[str, Any]) -> Dict[str, Any]:
        """Create PowerPoint presentation from analysis data"""
        request_id = request.get("request_id", "unknown")
        analysis_data = request.get("analysis_data", {})
        template = request.get("template", "comprehensive")
        output_path = request.get("output_path", f"reports/video_analysis_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
        
        try:
            # Ensure reports directory exists
            Path("reports").mkdir(exist_ok=True)
            
            # Create presentation
            prs = Presentation()
            
            # Slide 1: Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "SVA Video Analysis Report"
            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}\nVideo: {analysis_data.get('video_name', 'Unknown')}"
            
            # Slide 2: Executive Summary
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Executive Summary"
            tf = content.text_frame
            tf.text = "Key highlights from the video analysis:"
            
            summary_points = self._generate_summary_points(analysis_data)
            for point in summary_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
            
            # Slide 3: Audio Analysis (if available)
            if analysis_data.get("transcription"):
                slide = prs.slides.add_slide(bullet_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Audio Analysis"
                tf = content.text_frame
                
                trans_data = analysis_data["transcription"]
                tf.text = f"Language: {trans_data.get('language', 'Unknown')}"
                
                p = tf.add_paragraph()
                p.text = f"Confidence: {trans_data.get('confidence_score', 0):.1%}"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"Duration: {trans_data.get('duration', 0):.1f} seconds"
                p.level = 1
                
                # Add transcription excerpt
                text = trans_data.get('text', '')
                if text:
                    excerpt = text[:200] + "..." if len(text) > 200 else text
                    p = tf.add_paragraph()
                    p.text = f"Excerpt: {excerpt}"
                    p.level = 1
            
            # Slide 4: Visual Analysis (if available)
            if analysis_data.get("vision"):
                slide = prs.slides.add_slide(bullet_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Visual Content Analysis"
                tf = content.text_frame
                
                vision_data = analysis_data["vision"]
                
                # Objects
                if vision_data.get("objects_detected"):
                    object_counts = {}
                    for obj in vision_data["objects_detected"]:
                        class_name = obj["class_name"]
                        object_counts[class_name] = object_counts.get(class_name, 0) + 1
                    
                    tf.text = "Objects Detected:"
                    for obj_name, count in object_counts.items():
                        p = tf.add_paragraph()
                        p.text = f"{obj_name.title()}: {count} instances"
                        p.level = 1
                
                # Text content
                if vision_data.get("text_found"):
                    p = tf.add_paragraph()
                    p.text = "Text Content Found:"
                    p.level = 0
                    
                    for text_item in vision_data["text_found"][:3]:  # Limit to first 3
                        p = tf.add_paragraph()
                        p.text = f'"{text_item["text"]}"'
                        p.level = 1
            
            # Slide 5: Key Findings
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Key Findings & Recommendations"
            tf = content.text_frame
            tf.text = "Summary of important insights:"
            
            findings = self._generate_key_findings(analysis_data)
            for finding in findings:
                p = tf.add_paragraph()
                p.text = finding
                p.level = 1
            
            # Save presentation
            prs.save(output_path)
            
            # Get file size
            file_size = Path(output_path).stat().st_size
            
            logger.info(f"📊 PowerPoint presentation generated: {output_path}")
            
            return {
                "request_id": request_id,
                "status": "success",
                "data": {
                    "output_path": output_path,
                    "file_size": file_size,
                    "format": "PowerPoint",
                    "template_used": template,
                    "slides_created": len(prs.slides),
                    "sections_included": ["Title", "Executive Summary", "Audio Analysis", "Visual Analysis", "Key Findings"]
                }
            }
            
        except Exception as e:
            logger.error(f"❌ PowerPoint generation failed: {e}")
            return {
                "request_id": request_id,
                "status": "error",
                "error": str(e)
            }
    
    async def _create_multi_modal_report(self, request: Dict[str, Any]) -> Dict[str, Any]:
        """Create comprehensive multi-modal report combining all analysis"""
        request_id = request.get("request_id", "unknown")
        transcription_data = request.get("transcription_data", {})
        vision_data = request.get("vision_data", {})
        video_info = request.get("video_info", {})
        format_type = request.get("format", "pdf")  # pdf or pptx
        
        # Combine all analysis data
        combined_data = {
            "video_name": video_info.get("name", "Unknown Video"),
            "transcription": transcription_data,
            "vision": vision_data,
            "analysis_timestamp": datetime.now().isoformat()
        }
        
        if format_type.lower() == "pptx":
            return await self._create_powerpoint({
                "request_id": request_id,
                "analysis_data": combined_data,
                "template": "comprehensive"
            })
        else:
            return await self._generate_pdf_report({
                "request_id": request_id,
                "analysis_data": combined_data,
                "template": "comprehensive"
            })
    
    def _generate_executive_summary(self, analysis_data: Dict[str, Any]) -> str:
        """Generate executive summary from analysis data"""
        summary_parts = []
        
        # Video info
        video_name = analysis_data.get('video_name', 'the video')
        summary_parts.append(f"This report presents a comprehensive analysis of {video_name}.")
        
        # Audio analysis summary
        if analysis_data.get("transcription"):
            trans_data = analysis_data["transcription"]
            language = trans_data.get('language', 'unknown')
            confidence = trans_data.get('confidence_score', 0)
            word_count = trans_data.get('word_count', 0)
            
            summary_parts.append(f"The audio content was successfully transcribed with {confidence:.1%} confidence, detecting {language} language with approximately {word_count} words.")
        
        # Visual analysis summary
        if analysis_data.get("vision"):
            vision_data = analysis_data["vision"]
            objects = vision_data.get("objects_detected", [])
            text_items = vision_data.get("text_found", [])
            
            if objects:
                object_count = len(objects)
                unique_objects = len(set(obj["class_name"] for obj in objects))
                summary_parts.append(f"Visual analysis identified {object_count} objects across {unique_objects} different categories.")
            
            if text_items:
                summary_parts.append(f"Additionally, {len(text_items)} text elements were extracted from the video frames.")
        
        summary_parts.append("This multi-modal analysis provides valuable insights into both the audio and visual content of the video.")
        
        return " ".join(summary_parts)
    
    def _generate_summary_points(self, analysis_data: Dict[str, Any]) -> List[str]:
        """Generate bullet points for presentation summary"""
        points = []
        
        if analysis_data.get("transcription"):
            trans_data = analysis_data["transcription"]
            points.append(f"Language detected: {trans_data.get('language', 'Unknown')}")
            points.append(f"Transcription confidence: {trans_data.get('confidence_score', 0):.1%}")
        
        if analysis_data.get("vision"):
            vision_data = analysis_data["vision"]
            if vision_data.get("objects_detected"):
                object_types = set(obj["class_name"] for obj in vision_data["objects_detected"])
                points.append(f"Visual objects: {', '.join(list(object_types)[:3])}")
            
            if vision_data.get("text_found"):
                points.append(f"Text elements extracted: {len(vision_data['text_found'])}")
        
        points.append("Multi-modal analysis completed successfully")
        
        return points
    
    def _generate_key_findings(self, analysis_data: Dict[str, Any]) -> List[str]:
        """Generate key findings and recommendations"""
        findings = []
        
        # Audio findings
        if analysis_data.get("transcription"):
            trans_data = analysis_data["transcription"]
            if trans_data.get('confidence_score', 0) > 0.8:
                findings.append("High-quality audio transcription achieved with strong confidence scores")
            
            if trans_data.get('language') in ['ms', 'id']:
                findings.append("Successfully processed Malay/Indonesian language content")
        
        # Visual findings
        if analysis_data.get("vision"):
            vision_data = analysis_data["vision"]
            if vision_data.get("objects_detected"):
                findings.append("Clear visual objects identified, indicating good video quality")
            
            if vision_data.get("text_found"):
                findings.append("Text content successfully extracted from video frames")
        
        # General recommendations
        findings.append("Video content is suitable for automated analysis and processing")
        findings.append("Both audio and visual channels provide valuable information")
        
        return findings
    
    async def _get_report_templates(self) -> Dict[str, Any]:
        """Return available report templates"""
        return {
            "status": "success",
            "data": {
                "available_templates": self.report_templates,
                "supported_formats": self.supported_formats
            }
        }
    
    async def _get_capabilities(self) -> Dict[str, Any]:
        """Return server capabilities"""
        return {
            "status": "success",
            "data": {
                "server_name": self.server_name,
                "capabilities": self.capabilities,
                "supported_formats": self.supported_formats,
                "available_templates": list(self.report_templates.keys())
            }
        }

# Test functions
async def test_generation_mcp_server():
    """Test the generation MCP server"""
    print("🧪 Testing Generation MCP Server")
    print("=" * 50)
    
    # Initialize server
    server = GenerationMCPServer()
    init_result = await server.initialize()
    print(f"Initialization: {json.dumps(init_result, indent=2)}")
    
    if init_result["status"] != "success":
        print("❌ Server initialization failed!")
        return
    
    # Test 1: Get capabilities
    print("\n📋 Test 1: Get Capabilities")
    result = await server.process_request({"action": "get_capabilities"})
    print(f"Capabilities: {json.dumps(result, indent=2)}")
    
    # Test 2: Get templates
    print("\n📄 Test 2: Get Report Templates")
    result = await server.process_request({"action": "get_report_templates"})
    print(f"Templates: {json.dumps(result, indent=2)}")
    
    # Test 3: Create sample analysis data
    sample_analysis = {
        "video_name": "test_video.mp4",
        "transcription": {
            "text": "Kemerdekan, kemerdekan damak, asyik ulang, kemerdekan berhati, cedera usus perut rakyat berisi nasi, vitamin dan zat bisi.",
            "language": "ms",
            "confidence_score": 0.95,
            "duration": 43.56,
            "word_count": 29
        },
        "vision": {
            "objects_detected": [
                {"class_name": "person", "confidence": 0.76},
                {"class_name": "bottle", "confidence": 0.79},
                {"class_name": "bottle", "confidence": 0.82}
            ],
            "text_found": [
                {"text": "kemerdekaan", "confidence": 0.99},
                {"text": "vitamin", "confidence": 1.00},
                {"text": "nasi", "confidence": 0.95}
            ]
        }
    }
    
    # Test 3: Generate PDF Report
    print("\n📄 Test 3: Generate PDF Report")
    request = {
        "action": "generate_pdf_report",
        "analysis_data": sample_analysis,
        "template": "comprehensive"
    }
    result = await server.process_request(request)
    if result["status"] == "success":
        print(f"✅ PDF generated: {result['data']['output_path']}")
        print(f"File size: {result['data']['file_size']} bytes")
    else:
        print(f"❌ PDF generation failed: {result['error']}")
    
    # Test 4: Generate PowerPoint
    print("\n📊 Test 4: Generate PowerPoint Presentation")
    request = {
        "action": "create_powerpoint",
        "analysis_data": sample_analysis,
        "template": "comprehensive"
    }
    result = await server.process_request(request)
    if result["status"] == "success":
        print(f"✅ PowerPoint generated: {result['data']['output_path']}")
        print(f"Slides created: {result['data']['slides_created']}")
    else:
        print(f"❌ PowerPoint generation failed: {result['error']}")
    
    # Test 5: Multi-modal report
    print("\n🔄 Test 5: Create Multi-Modal Report")
    request = {
        "action": "create_multi_modal_report",
        "transcription_data": sample_analysis["transcription"],
        "vision_data": sample_analysis["vision"],
        "video_info": {"name": "test_video.mp4"},
        "format": "pdf"
    }
    result = await server.process_request(request)
    if result["status"] == "success":
        print(f"✅ Multi-modal report generated: {result['data']['output_path']}")
    else:
        print(f"❌ Multi-modal report failed: {result['error']}")
    
    print("\n✅ Generation MCP Server testing completed!")

if __name__ == "__main__":
    asyncio.run(test_generation_mcp_server())